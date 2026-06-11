from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Title Slide (Layout 0)
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Adversarial Geometric Distillation"
subtitle.text = "Spatial Grounding for Hallucination Mitigation in Text-to-3D Generation\n\nGiman Dissanayaka | 215519C\nUniversity of Moratuwa"

# Helper function to add slides
def add_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, pt in enumerate(bullet_points):
        if i == 0:
            tf.text = pt
        else:
            p = tf.add_paragraph()
            if pt.startswith("  -"):
                p.text = pt.strip("- ")
                p.level = 1
            else:
                p.text = pt
                p.level = 0
    return slide

# Introduction
add_slide("What is Text-to-3D Generation?", [
    "AI generates 3D models from text prompts (e.g., 'a knight in armour')",
    "Powered by Score Distillation Sampling (SDS) + pre-trained 2D diffusion priors",
    "Key frameworks: DreamFusion, Magic3D, Fantasia3D"
])

add_slide("The Problem: 3D Hallucinations", [
    "What are hallucinations?",
    "  - Janus Faces — same face on front and back",
    "  - Floating Blobs — disconnected geometry fragments",
    "  - Geometric Spikes — sharp protrusions on the surface",
    "  - Topology Errors — holes, handles, non-manifold patches",
    "Why do they occur?",
    "  - SDS gradient is geometry-blind — defined entirely in 2D pixel space"
])

# Objectives
add_slide("Research Objectives", [
    "1. Detect geometric hallucinations using 3D-native spectral and topological metrics",
    "2. Spatially ground detections to individual mesh vertices without language intermediary",
    "3. Apply a bounded, selective geometric refinement using grounded vertex weights",
    "4. Validate the closed-loop pipeline on a real 3D mesh dataset",
    "Core Hypothesis: A spatially-aware geometric feedback mechanism grounded in spectral graph theory and algebraic topology is both necessary and sufficient to eliminate hallucinations without re-prompting a diffusion model"
])

# Literature Review
add_slide("Evolution of Text-to-3D", [
    "2017: Point-E / Point Set Gen. - First learned 3D point cloud synthesis",
    "2022: Dream Fields / Latent-NeRF - CLIP-guided NeRF optimisation",
    "2022: DreamFusion (SDS) - 2D diffusion prior → 3D via score distillation",
    "2023: Magic3D, Fantasia3D - Coarse-to-fine, disentangled geometry/appearance",
    "2023: ProlificDreamer (VSD) - Variational score distillation",
    "2023: Hallo3D - First LMM-based hallucination detection + language feedback",
    "2024: MeshGPT, 3D Gaussian Splatting"
])

add_slide("How SDS Works & Limitations", [
    "SDS frames 3D generation as distilling a 2D diffusion prior",
    "Gradient flows from 2D pixel space → 3D representation",
    "Limitation: No geometric constraint exists in this objective"
])

add_slide("Existing Mitigation Approaches", [
    "Perp-Neg: Modifies SDS gradient direction (Only reduces Janus, cannot fix geometry)",
    "Entropic SDS: Score-space debiasing (Geometry-blind — no mesh inspection)",
    "Hallo3D: LMM inspects renders → language prompt (Semantic bottleneck: spatial info lost in text)",
    "Laplacian Mesh Editing (Sorkine): Energy minimisation on vertex positions (Requires explicit spatial constraints)"
])

add_slide("What is the Semantic Bottleneck?", [
    "Hallo3D detects: 'Janus face at mesh back region'",
    "Converts to language: negative_prompt = 'no Janus face'",
    "Feeds back to diffusion model → diffuse, global correction",
    "Result:",
    "  - Vertex coordinates lost",
    "  - Topology info lost",
    "  - Spatial precision lost"
])

# Problem Definition
add_slide("Formal Research Gap", [
    "No existing framework implements a closed-loop feedback mechanism that detects geometric hallucinations using spatially complete 3D metrics, grounds detections directly to individual mesh vertices without natural-language conversion, and applies a spatially selective, theoretically bounded geometric correction.",
    "Three specific gaps:",
    "  - SDS objective has no geometric energy term — geometry-blind by design",
    "  - No vertex-level correction signal derived from 3D analysis exists",
    "  - Language channel (Hallo3D) discards spatial precision at the feedback stage"
])

# Technology Adopted
add_slide("Spectral Graph Theory as a Detection Tool", [
    "Mesh = Graph G = (V, E) → Laplacian: L = D − A",
    "Eigenvalue Variance: High variance → spiky, irregular geometry",
    "Fiedler Value (λ2): Near zero → mesh close to disconnected (floating artefacts)",
    "Computed in O(|V|+|E|) — scalable to 500K+ vertex meshes"
])

add_slide("Algebraic Topology as a Global Integrity Check", [
    "Euler Characteristic: χ = V − E + F → Genus: g = (2 − χ)/2",
    "Betti Numbers: β0 = connected components, β1 = loops, β2 = voids",
    "Metrics:",
    "  - β0 (Components): 1 (Clean) vs > 1 (floating blobs)",
    "  - Genus g: 0 (Clean) vs > 0 (unintended handles)",
    "  - Fiedler λ2: > 0 (Clean) vs ≈ 0 (near-disconnected)"
])

add_slide("GNN Critic + Laplacian Refinement", [
    "Graph Convolutional Critic: 2-layer GCNConv maps vertex features → anomaly score per vertex",
    "  - Input: [degree(vi), mean_edge_length(vi)]",
    "  - Output: anomaly score si ∈ [0, 1]",
    "Laplacian Mesh Editing (Sorkine): Energy minimisation for bounded deformation",
    "  - E(x) = λ ||Lv||^2 + (1-λ)||v - v(0)||^2"
])

# Novel Approach
add_slide("AGD Framework Overview", [
    "AGD = closed-loop refinement layer that attaches to any text-to-3D generator",
    "Three modules:",
    "  1. Preprocessing — load, validate, normalise mesh",
    "  2. ML Engine — discriminator + critic + LMM detector",
    "  3. Extended Module — geometric grounding + refinement optimizer"
])

add_slide("The Core Innovation: Geometric Grounding", [
    "Problem: How to convert 'back view severity = 0.7' → which vertices to fix?",
    "Solution: Projection-based grounding — no language required",
    "  1. Project all vertices onto view direction dv",
    "  2. Select top quantile by projection (closest to camera)",
    "  3. Filter: normal alignment ni · dv ≥ threshold",
    "  4. Assign: wi = max(wi, severity)",
    "Result: Per-vertex weight map wi ∈ [0,1] with spatial precision at vertex level"
])

add_slide("The Refinement Loop", [
    "Weighted energy minimisation applied per-vertex:",
    "v_i ← v_i - η · (λ ∇L_geo + (1-λ)∇D) · w_i",
    "Components:",
    "  - L_geo = Laplacian smoothing (removes spikes)",
    "  - D = Distillation anchor (prevents over-smoothing clean regions)",
    "  - w_i = grounded weight (concentrates correction on hallucinated vertices only)"
])

# Current Implementation
add_slide("EXP_03: CLIP-Based Detection Baseline", [
    "Goal: Validate training-free hallucination detection using CLIP + geometric metrics",
    "5 Detectors on 24 rendered views (8 azimuth × 3 elevation, 512×512):",
    "  - CLIP Semantic Deviation (35%): Janus faces, semantic drift",
    "  - Depth Discontinuity (25%): Floating geometry",
    "  - Normal Inconsistency (20%): Inverted surfaces",
    "  - Silhouette Asymmetry (15%): Asymmetric Janus geometry",
    "  - Edge Density Variance (5%): View structural inconsistency"
])

add_slide("EXP_03: Dragon Mesh Results", [
    "Test mesh: Stanford Dragon (~430K vertices)",
    "Global hallucination score: 0.0856 (Clean)",
    "No Janus, floater, or normal-flip detected",
    "Dominant loss: L_depth = 0.017107 (expected for complex surface)",
    "View weights: Uniform ~0.042 (no single view dominates)"
])

add_slide("EXP_04: AGD Pipeline — What's Built", [
    "Fully Implemented:",
    "  - Discriminator: topology (χ, genus, β0) + spectral + geometry quality",
    "  - Heuristic Critic: degree z-score + edge-length z-score per vertex",
    "  - Renderer: Fibonacci-sphere, 6–100+ views, 5 buffers",
    "  - Grounding Module: view-to-vertex projection + max-fusion",
    "  - Refinement Optimizer: Laplacian + anchor, weighted per-vertex",
    "  - LMM Detector: Local LLaVA v1.5/1.6, JSON severity output",
    "Planned:",
    "  - SDS/Threestudio integration | Differentiable topology losses | Objaverse benchmarks"
])

add_slide("AGD Pipeline — Key Results", [
    "Pipeline executes end-to-end on 21-mesh test dataset",
    "Before/after hallucination score comparison printed per mesh",
    "10-metric geometry error report: CD, HD95, NCS, SSIM, Angular Error, Roughness RMS, Edge-IoU",
    "Example CLI:",
    "  python agd_pipeline.py --input-file 3d_samples/janus.obj \\",
    "  --render-views --num-views 36 --use-vlm"
])

add_slide("Test Dataset", [
    "21 meshes including clean and hallucinated samples",
    "Categories:",
    "  - Clean references: bunny.ply, dragon.ply, happy.ply, horse.ply",
    "  - AI-generated: Meshy_AI_Infernal_Ironclad.obj",
    "  - Known hallucinations: janus.obj, Mr_Bean.obj",
    "  - Character/organic: Bearded_guy.ply, Rigged_Hand.obj, hand.ply",
    "  - Complex geometry: blade.ply, venusm.obj, elepham.obj"
])

add_slide("EXP_03 → EXP_04: Key Upgrades", [
    "Views:",
    "  - EXP_03: 24 (fixed grid) | EXP_04: 6–100+ (Fibonacci-sphere)",
    "Detection:",
    "  - EXP_03: 5 image detectors | EXP_04: Graph-theoretic + LMM + critic",
    "Correction:",
    "  - EXP_03: Loss scalar only | EXP_04: Per-vertex weight map",
    "Spatial precision:",
    "  - EXP_03: View-level | EXP_04: Vertex-level",
    "Bottleneck:",
    "  - EXP_03: Loss scalar | EXP_04: Eliminated entirely"
])

# Design
add_slide("System Architecture (Three Modules)", [
    "Preprocessing Module (agd_pipeline.py): Load → Validate → Normalise → Adjacency graph",
    "ML Engine (discriminator.py, critic.py, renderer.py, lmm_detector.py): Detect hallucinations using global + local + visual signals",
    "Extended Module (grounding.py, optimizer.py): Ground → Refine → Evaluate"
])

add_slide("Data Flow Design", [
    "Key interaction patterns:",
    "  1. Sequential execution with optional LMM/rendering branches",
    "  2. Signal fusion at grounding: max(critic_weight, LMM_bias) per vertex",
    "  3. Discriminator called before and after refinement (symmetric)",
    "  4. 10-metric geometry report from geometry_metrics.py"
])

add_slide("Integration Roadmap", [
    "Level 1 (Current) - Post-processing: any .obj/.ply/.stl → AGD pipeline (Working)",
    "Level 2 (Planned) - In-loop: Add E(x) to SDS loss in Threestudio (Planned)",
    "Level 3 (Planned) - Benchmark: Objaverse/GSO evaluation + comparison tables (Planned)"
])

# References
add_slide("Key References", [
    "1. Poole et al. — DreamFusion: Text-to-3D Using 2D Diffusion",
    "2. Wang et al. — ProlificDreamer: VSD for Text-to-3D",
    "3. Chen et al. — Fantasia3D: Disentangling Geometry and Appearance",
    "4. Cui et al. — Hallo3D: Dynamic Portrait Animation",
    "5. Armandpour et al. — Perp-Neg: Alleviating Janus Problem",
    "6. Sorkine — Laplacian Mesh Processing",
    "7. Liu et al. — Visual Instruction Tuning (LLaVA)"
])

prs.save("AGD_Presentation.pptx")
